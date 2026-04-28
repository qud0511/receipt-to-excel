import base64
import io

from PIL import Image
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

from . import ProcessedInput


def process_presentation(file_bytes: bytes, source_name: str) -> list[ProcessedInput]:
    prs = Presentation(io.BytesIO(file_bytes))
    results: list[ProcessedInput] = []

    for slide_num, slide in enumerate(prs.slides):
        embedded_images: list[Image.Image] = []
        text_parts: list[str] = []

        for shape in slide.shapes:
            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                img = Image.open(io.BytesIO(shape.image.blob))
                embedded_images.append(img)
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    t = para.text.strip()
                    if t:
                        text_parts.append(t)

        if embedded_images:
            img = embedded_images[0].convert("RGB")
            buf = io.BytesIO()
            img.save(buf, format="PNG")
            b64 = base64.b64encode(buf.getvalue()).decode()
            results.append(ProcessedInput(
                source_name=source_name,
                source_page=slide_num,
                image_b64=b64,
                text=None,
                pil_image=img.copy(),
            ))
        elif text_parts:
            results.append(ProcessedInput(
                source_name=source_name,
                source_page=slide_num,
                image_b64=None,
                text="\n".join(text_parts),
                pil_image=None,
            ))

    return results
