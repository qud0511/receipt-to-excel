import io
import pytest
from PIL import Image as PilImage
from app.services.preprocessor import ProcessedInput, route_file


def test_processed_input_has_docling_text_field():
    pi = ProcessedInput(
        source_name="test.jpg",
        source_page=0,
        docling_text="업체명: 스타벅스\n금액: 5500",
        pil_image=None,
    )
    assert pi.docling_text == "업체명: 스타벅스\n금액: 5500"
    assert not hasattr(pi, "image_b64")
    assert not hasattr(pi, "text")


def test_unsupported_extension_raises():
    with pytest.raises(ValueError, match="Unsupported"):
        route_file(b"data", "document.txt")


# ── fixtures ──────────────────────────────────────────────────────────────────

@pytest.fixture
def white_jpg_bytes() -> bytes:
    img = PilImage.new("RGB", (100, 80), color=(255, 255, 255))
    buf = io.BytesIO()
    img.save(buf, format="JPEG")
    return buf.getvalue()


@pytest.fixture
def white_png_bytes() -> bytes:
    img = PilImage.new("RGB", (60, 40), color=(200, 200, 200))
    buf = io.BytesIO()
    img.save(buf, format="PNG")
    return buf.getvalue()


@pytest.fixture
def two_page_pdf_bytes() -> bytes:
    import fitz
    doc = fitz.open()
    for i in range(2):
        page = doc.new_page(width=200, height=150)
        page.insert_text((50, 75), f"Page {i + 1}")
    buf = io.BytesIO()
    doc.save(buf)
    doc.close()
    return buf.getvalue()


@pytest.fixture
def simple_xlsx_bytes() -> bytes:
    from openpyxl import Workbook
    wb = Workbook()
    ws = wb.active
    ws["A1"] = "날짜"
    ws["B1"] = "금액"
    ws["A2"] = "2024-01-15"
    ws["B2"] = 5500
    buf = io.BytesIO()
    wb.save(buf)
    return buf.getvalue()


@pytest.fixture
def two_slide_pptx_bytes() -> bytes:
    from pptx import Presentation as Prs
    from pptx.util import Inches
    prs = Prs()
    for i in range(2):
        slide = prs.slides.add_slide(prs.slide_layouts[5])
        tb = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(4), Inches(2))
        tb.text_frame.text = f"영수증 {i + 1}\n금액: {(i + 1) * 1000}원"
    buf = io.BytesIO()
    prs.save(buf)
    return buf.getvalue()


# ── image tests ───────────────────────────────────────────────────────────────

def test_jpg_returns_single_input(white_jpg_bytes):
    results = route_file(white_jpg_bytes, "receipt.jpg")
    assert len(results) == 1
    r = results[0]
    assert r.source_name == "receipt.jpg"
    assert r.source_page == 0
    assert r.docling_text is not None
    assert r.pil_image is not None


def test_png_has_pil_image(white_png_bytes):
    results = route_file(white_png_bytes, "scan.png")
    r = results[0]
    assert r.pil_image is not None
    assert r.pil_image.size[0] > 0


# ── PDF tests ─────────────────────────────────────────────────────────────────

def test_pdf_one_input_per_page(two_page_pdf_bytes):
    results = route_file(two_page_pdf_bytes, "invoice.pdf")
    assert len(results) == 2


def test_pdf_page_metadata(two_page_pdf_bytes):
    results = route_file(two_page_pdf_bytes, "invoice.pdf")
    assert results[0].source_page == 0
    assert results[1].source_page == 1
    assert results[0].docling_text
    assert results[0].pil_image is not None


# ── XLSX tests ────────────────────────────────────────────────────────────────

def test_xlsx_returns_single_text_input(simple_xlsx_bytes):
    results = route_file(simple_xlsx_bytes, "data.xlsx")
    assert len(results) == 1
    r = results[0]
    assert r.pil_image is None
    assert r.docling_text
    assert "5500" in r.docling_text


# ── PPTX tests ────────────────────────────────────────────────────────────────

def test_pptx_one_input_per_slide(two_slide_pptx_bytes):
    results = route_file(two_slide_pptx_bytes, "slides.pptx")
    assert len(results) == 2


def test_pptx_text_extracted(two_slide_pptx_bytes):
    results = route_file(two_slide_pptx_bytes, "slides.pptx")
    combined = results[0].docling_text + results[1].docling_text
    assert "영수증" in combined
    assert results[0].source_page == 0
    assert results[1].source_page == 1
